from pptx import Presentation
from pptx.util import Inches

# Data
questions_answers = [
    ("What is the total number of impressions across all ads?", "213434828"),
    ("What is the total number of clicks across all ads?", "38165"),
    (
        "What is the average click-through rate (CTR) across all ads?",
        "0.00017881336592357833",
    ),
    ("What is the total cost spent across all ads?", "58705.229958205004"),
    ("What is the average cost per click (CPC) across all ads?", "0.042827"),
    ("What is the total number of conversions across all ads?", "3264"),
    ("What is the total number of approved conversions across all ads?", "1079"),
    (
        "What is the average conversion rate (CVR) across all ads?",
        "0.08552338530066815",
    ),
    ("What is the average cost per acquisition (CPA) across all ads?", "54.41"),
    (
        "What is the average approved conversion rate (ACR) across all ads?",
        "0.028271976942224553",
    ),
]

# Create a new presentation
prs = Presentation()

# Iterate through questions and answers
for idx, (question, answer) in enumerate(questions_answers):
    # Add a new slide
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    # Add question as title
    title = slide.shapes.title
    title.text = f"Question {idx}: {question}"

    # Add answer as content
    content = slide.placeholders[1]
    content.text = f"Answer: {answer}"

    # Add sentiment/commentary placeholder
    left = Inches(1)
    top = Inches(3)
    width = Inches(8)
    height = Inches(2)
    shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = shape.text_frame
    p = text_frame.add_paragraph()

    # Add sentiment/commentary
    p.text = "Commentary: This is a placeholder for sentiment analysis or commentary."

# Save the presentation
prs.save("questions_answers.pptx")
